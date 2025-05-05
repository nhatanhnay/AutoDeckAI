from pptx import Presentation
from pptx.util import Inches
import tempfile

def generate_ppt(word_text, analysis_text, plot_path):
    prs = Presentation()
    layout = prs.slide_layouts[1]

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Tóm tắt từ Word"
    slide1.placeholders[1].text = word_text if word_text.strip() else "(Không có nội dung)"

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Phân tích dữ liệu"
    slide2.placeholders[1].text = analysis_text

    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Biểu đồ minh họa"
    slide3.shapes.add_picture(plot_path, Inches(1), Inches(2), width=Inches(8))

    output_path = tempfile.mktemp(suffix=".pptx")
    prs.save(output_path)
    return output_path